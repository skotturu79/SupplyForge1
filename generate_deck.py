from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

DARK=RGBColor(0x1F,0x29,0x37)
EM=RGBColor(0x10,0xB9,0x81)
EMD=RGBColor(0x05,0x96,0x69)
WHITE=RGBColor(0xFF,0xFF,0xFF)
OFFWH=RGBColor(0xF3,0xF4,0xF6)
GRAY=RGBColor(0x6B,0x72,0x80)
LG=RGBColor(0xE5,0xE7,0xEB)
BLUE=RGBColor(0x3B,0x82,0xF6)
AMBER=RGBColor(0xF5,0x9E,0x0B)
PURPLE=RGBColor(0x8B,0x5C,0xF6)
CYAN=RGBColor(0x06,0xB6,0xD4)
INDIGO=RGBColor(0x63,0x66,0xF1)
RED=RGBColor(0xEF,0x44,0x44)
DKCARD=RGBColor(0x2D,0x3A,0x4A)
SW=Inches(13.33)
SH=Inches(7.5)

def rect(sl,x,y,w,h,fill=None,line=None,lw=None):
    sh=sl.shapes.add_shape(1,Inches(x),Inches(y),Inches(w),Inches(h))
    if fill:
        sh.fill.solid(); sh.fill.fore_color.rgb=fill
    else:
        sh.fill.background()
    if line:
        sh.line.color.rgb=line; sh.line.width=lw if lw else Pt(1)
    else:
        sh.line.fill.background()
    return sh

def tb(sl,x,y,w,h,text,sz=12,bold=False,col=WHITE,al=PP_ALIGN.LEFT,it=False):
    t=sl.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h))
    tf=t.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.alignment=al
    r=p.add_run(); r.text=text
    r.font.size=Pt(sz); r.font.bold=bold; r.font.italic=it; r.font.color.rgb=col
    return t

def oval(sl,cx,cy,r,fill):
    sh=sl.shapes.add_shape(9,Inches(cx-r),Inches(cy-r),Inches(r*2),Inches(r*2))
    sh.fill.solid(); sh.fill.fore_color.rgb=fill; sh.line.fill.background()
    return sh

def pgn(sl,n):
    tb(sl,12.6,7.15,0.6,0.25,str(n),sz=9,col=GRAY,al=PP_ALIGN.RIGHT)

def bcode(sl,x,y,w,h,col=DARK):
    pat=[3,1,2,1,3,2,1,2,1,3,1,1,2,3,1,2,1,1,3,2,1,2,3,1,2,1,2,3,1,1]
    tot=sum(pat); unit=w/tot; cx=x
    for i,u in enumerate(pat):
        if i%2==0:
            rect(sl,cx,y,max(unit*u,0.005),h,fill=col)
        cx+=unit*u

def qr(sl,x,y,cs):
    m=[[1,1,1,1,1,1,1],[1,0,0,0,0,0,1],[1,0,1,1,1,0,1],
       [1,0,1,0,1,0,1],[1,0,1,1,1,0,1],[1,0,0,0,0,0,1],[1,1,1,1,1,1,1]]
    d=[[0,0,0,0,0,0,0,0],[0,1,0,1,1,0,1,0],[0,0,1,0,0,1,0,0],
       [0,1,1,0,1,0,1,0],[0,0,0,1,0,1,0,0]]
    for r2 in range(7):
        for c2 in range(7):
            if m[r2][c2]: rect(sl,x+c2*cs,y+r2*cs,cs*0.88,cs*0.88,fill=DARK)
    ox=x+8*cs
    for r2 in range(7):
        for c2 in range(7):
            if m[r2][c2]: rect(sl,ox+c2*cs,y+r2*cs,cs*0.88,cs*0.88,fill=DARK)
    oy=y+8*cs
    for r2 in range(7):
        for c2 in range(7):
            if m[r2][c2]: rect(sl,x+c2*cs,oy+r2*cs,cs*0.88,cs*0.88,fill=DARK)
    for r2 in range(len(d)):
        for c2 in range(len(d[0])):
            if d[r2][c2]: rect(sl,x+(c2+7)*cs,y+(r2+7)*cs,cs*0.82,cs*0.82,fill=DARK)

def dm(sl,x,y,cs):
    m=[[1,1,1,1,1,1,1,1],[1,0,1,0,1,1,0,0],[1,1,0,1,0,0,1,0],[1,0,1,1,1,0,0,0],
       [1,1,0,0,0,1,1,0],[1,0,1,0,1,0,0,0],[1,1,0,1,0,1,1,0],[1,0,1,0,1,0,1,1]]
    for r2 in range(8):
        for c2 in range(8):
            if m[r2][c2]: rect(sl,x+c2*cs,y+r2*cs,cs*0.88,cs*0.88,fill=DARK)

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

