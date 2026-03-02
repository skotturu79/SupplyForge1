"""
SupplyForge — Summary Presentation Generator
Generates a professional 14-slide PowerPoint deck.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Brand colours ──────────────────────────────────────────────────
NAVY       = RGBColor(0x35, 0x4A, 0x5E)   # #354A5E  (header/bg)
DARK_NAVY  = RGBColor(0x1B, 0x2B, 0x3A)   # #1B2B3A
GREEN      = RGBColor(0x10, 0x7E, 0x3E)   # #107E3E  (accent)
LIGHT_GREEN= RGBColor(0x34, 0xD3, 0x99)   # #34D399
BLUE       = RGBColor(0x00, 0x70, 0xF2)   # #0070F2  (info)
AMBER      = RGBColor(0xE9, 0x73, 0x0C)   # #E9730C
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF0, 0xF2, 0xF4)   # #F0F2F4
MID_GRAY   = RGBColor(0x6A, 0x6D, 0x70)
DARK_GRAY  = RGBColor(0x32, 0x36, 0x3A)
BORDER     = RGBColor(0xED, 0xED, 0xED)
PURPLE     = RGBColor(0x6B, 0x3F, 0xA0)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # completely blank

# ── Helper functions ───────────────────────────────────────────────

def add_rect(slide, left, top, width, height, fill_rgb=None, line_rgb=None, line_width=Pt(0)):
    from pptx.util import Pt
    shape = slide.shapes.add_shape(1, left, top, width, height)  # 1 = MSO_SHAPE_TYPE.RECTANGLE
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if line_rgb:
        shape.line.color.rgb = line_rgb
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height,
             font_size=Pt(12), bold=False, color=DARK_GRAY,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf    = txBox.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = font_size
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txBox


def add_para(tf, text, font_size=Pt(11), bold=False, color=DARK_GRAY,
             align=PP_ALIGN.LEFT, space_before=Pt(4), italic=False, level=0):
    from pptx.util import Pt
    p = tf.add_paragraph()
    p.alignment  = align
    p.space_before = space_before
    p.level = level
    run = p.add_run()
    run.text = text
    run.font.size   = font_size
    run.font.bold   = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return p


def header_bar(slide, title, subtitle=None):
    """Dark navy header bar at top of slide."""
    add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.25), fill_rgb=NAVY)
    add_text(slide, title, Inches(0.4), Inches(0.15), Inches(10), Inches(0.65),
             font_size=Pt(24), bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle, Inches(0.4), Inches(0.75), Inches(10), Inches(0.4),
                 font_size=Pt(12), color=RGBColor(0xB0, 0xC4, 0xD8), italic=True)
    # Green accent line
    add_rect(slide, Inches(0), Inches(1.25), prs.slide_width, Inches(0.05), fill_rgb=GREEN)


def card(slide, left, top, width, height, title=None, title_color=NAVY,
         bg=WHITE, border=BORDER, accent_color=None):
    """Card with optional top accent stripe."""
    add_rect(slide, left, top, width, height, fill_rgb=bg, line_rgb=border, line_width=Pt(0.75))
    if accent_color:
        add_rect(slide, left, top, width, Inches(0.07), fill_rgb=accent_color)
    if title:
        add_text(slide, title, left + Inches(0.15), top + Inches(0.1),
                 width - Inches(0.3), Inches(0.35),
                 font_size=Pt(10), bold=True, color=title_color)
    return top + (Inches(0.45) if title else Inches(0.15))


def bullet(slide, txBox_or_tf, items, icon="▸", font_size=Pt(11),
           color=DARK_GRAY, indent=Inches(0.15), space=Pt(5)):
    if hasattr(txBox_or_tf, 'text_frame'):
        tf = txBox_or_tf.text_frame
    else:
        tf = txBox_or_tf
    for item in items:
        add_para(tf, f"{icon}  {item}", font_size=font_size, color=color, space_before=space)


def sf_logo(slide, left, top, size=Inches(0.55)):
    """SF logo badge."""
    add_rect(slide, left, top, size, size, fill_rgb=BLUE)
    add_text(slide, "SF", left, top, size, size,
             font_size=Pt(16), bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)

# Full background gradient simulation (2 rects)
add_rect(slide, Inches(0), Inches(0), prs.slide_width, prs.slide_height, fill_rgb=DARK_NAVY)
add_rect(slide, Inches(0), Inches(4.5), prs.slide_width, Inches(3), fill_rgb=RGBColor(0x10, 0x35, 0x22))

# Green accent stripe
add_rect(slide, Inches(0), Inches(3.8), Inches(13.33), Inches(0.04), fill_rgb=GREEN)

# SF logo
sf_logo(slide, Inches(0.6), Inches(1.2), size=Inches(0.9))

# Title
add_text(slide, "SupplyForge", Inches(0.55), Inches(2.25), Inches(10), Inches(1.1),
         font_size=Pt(54), bold=True, color=WHITE)
add_text(slide, "Enterprise B2B Supply Chain Management Platform",
         Inches(0.55), Inches(3.3), Inches(10), Inches(0.55),
         font_size=Pt(18), color=LIGHT_GREEN, italic=True)

# Tag pills
pill_data = [("Multi-Tenant SaaS", BLUE), ("AI-Powered", GREEN), ("Enterprise-Grade", AMBER)]
x = Inches(0.55)
for label, color in pill_data:
    w = Inches(1.9)
    add_rect(slide, x, Inches(4.1), w, Inches(0.38),
             fill_rgb=RGBColor(color[0]//3, color[1]//3, color[2]//3),
             line_rgb=color, line_width=Pt(1.2))
    add_text(slide, label, x, Inches(4.1), w, Inches(0.38),
             font_size=Pt(10), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    x += Inches(2.1)

# Bottom meta
add_text(slide, "Version 1.0  ·  February 2026  ·  Confidential",
         Inches(0.55), Inches(6.8), Inches(9), Inches(0.4),
         font_size=Pt(10), color=MID_GRAY, italic=True)


# ══════════════════════════════════════════════════════════════════
# SLIDE 2 — EXECUTIVE SUMMARY
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, Inches(0), Inches(0), prs.slide_width, prs.slide_height, fill_rgb=LIGHT_GRAY)
header_bar(slide, "Executive Summary", "What is SupplyForge?")

# Left intro block
add_rect(slide, Inches(0.35), Inches(1.5), Inches(5.5), Inches(5.5), fill_rgb=WHITE,
         line_rgb=BORDER, line_width=Pt(0.75))
add_rect(slide, Inches(0.35), Inches(1.5), Inches(5.5), Inches(0.08), fill_rgb=GREEN)

tx = slide.shapes.add_textbox(Inches(0.55), Inches(1.65), Inches(5.1), Inches(5.0))
tf = tx.text_frame
tf.word_wrap = True
add_para(tf, "SupplyForge", font_size=Pt(20), bold=True, color=NAVY, space_before=Pt(2))
add_para(tf, "SupplyForge is a production-ready, multi-tenant B2B supply chain SaaS platform "
         "connecting businesses and their supplier/vendor network on a single unified platform.",
         font_size=Pt(11), color=MID_GRAY, space_before=Pt(8))

add_para(tf, "Core Value Propositions", font_size=Pt(12), bold=True, color=NAVY, space_before=Pt(16))
props = [
    "End-to-end document exchange (PO → Invoice → ASN → Delivery)",
    "Real-time shipment tracking across all major carriers",
    "AI-powered invoice matching via Claude Anthropic",
    "EDI support — X12 and EDIFACT for legacy ERP integration",
    "Enterprise security — AES-256-GCM, RSA-PSS, MFA",
    "Native multi-tenancy via PostgreSQL Row-Level Security",
    "Webhook + Data Feed ecosystem for real-time integrations",
]
for p in props:
    add_para(tf, f"✓  {p}", font_size=Pt(10.5), color=DARK_GRAY, space_before=Pt(5))

# Right KPI tiles
kpis = [
    ("3",       "Portal types",       GREEN,   Inches(6.2),  Inches(1.5)),
    ("16+",     "Vendor pages",       BLUE,    Inches(8.7),  Inches(1.5)),
    ("15+",     "API modules",        AMBER,   Inches(11.2), Inches(1.5)),
    ("8",       "MCP AI tools",       PURPLE,  Inches(6.2),  Inches(3.2)),
    ("4",       "Shared packages",    NAVY,    Inches(8.7),  Inches(3.2)),
    ("6+",      "Infrastructure svcs",GREEN,   Inches(11.2), Inches(3.2)),
    ("AES-256", "Encryption at rest", BLUE,    Inches(6.2),  Inches(4.9)),
    ("JWT RS256","Auth standard",      AMBER,   Inches(8.7),  Inches(4.9)),
    ("EDI",     "X12 + EDIFACT",      PURPLE,  Inches(11.2), Inches(4.9)),
]
for val, label, color, x, y in kpis:
    w, h = Inches(2.2), Inches(1.5)
    add_rect(slide, x, y, w, h, fill_rgb=WHITE, line_rgb=BORDER, line_width=Pt(0.75))
    add_rect(slide, x, y, w, Inches(0.07), fill_rgb=color)
    add_text(slide, val, x, y + Inches(0.15), w, Inches(0.75),
             font_size=Pt(26), bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(slide, label, x, y + Inches(0.85), w, Inches(0.55),
             font_size=Pt(9.5), color=MID_GRAY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════
# SLIDE 3 — ARCHITECTURE OVERVIEW
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, Inches(0), Inches(0), prs.slide_width, prs.slide_height, fill_rgb=LIGHT_GRAY)
header_bar(slide, "Architecture Overview", "Monorepo · Turborepo + pnpm workspaces")

layers = [
    ("PRESENTATION LAYER  —  Next.js 14 App Router (port 3000)",
     ["Admin Portal", "Business Portal", "Vendor Portal", "Public / Auth Pages"],
     NAVY, Inches(1.5)),
    ("API GATEWAY  —  Kong  (rate limiting · auth · request transformation)",
     [],
     AMBER, Inches(2.65)),
    ("APPLICATION LAYER  —  NestJS 10 + Fastify  (port 3001)",
     ["Auth · Documents · Partners · Vendors · Carriers · Tracking · Analytics",
      "API Keys · Webhooks · Notifications · Search · Storage · Tenants · AI"],
     GREEN, Inches(3.2)),
    ("DATA LAYER  —  PostgreSQL 16 · Redis 7 · Elasticsearch 8 · MinIO",
     ["Row-Level Security  ·  BullMQ job queues  ·  Full-text search  ·  Object storage"],
     BLUE, Inches(4.6)),
    ("MESSAGE BUS  —  RabbitMQ 3  (async events · webhook delivery · EDI routing)",
     [],
     PURPLE, Inches(5.55)),
]

for title, items, color, top in layers:
    h = Inches(0.55) if not items else Inches(0.55 + 0.32 * len(items))
    add_rect(slide, Inches(0.35), top, Inches(12.6), h, fill_rgb=WHITE,
             line_rgb=color, line_width=Pt(1.5))
    add_rect(slide, Inches(0.35), top, Inches(0.08), h, fill_rgb=color)
    add_text(slide, title, Inches(0.55), top + Inches(0.1), Inches(12.0), Inches(0.35),
             font_size=Pt(10), bold=True, color=color)
    y_off = Inches(0.38)
    for item in items:
        add_text(slide, f"  {item}", Inches(0.55), top + y_off, Inches(12.0), Inches(0.28),
                 font_size=Pt(9.5), color=MID_GRAY)
        y_off += Inches(0.28)

# Side label
add_text(slide, "MCP Server\n(Claude AI\nIntegration)",
         Inches(11.8), Inches(1.5), Inches(1.4), Inches(1.0),
         font_size=Pt(8.5), color=NAVY, bold=True, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════
# SLIDE 4 — THREE PORTALS
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, Inches(0), Inches(0), prs.slide_width, prs.slide_height, fill_rgb=LIGHT_GRAY)
header_bar(slide, "Three Portals", "One platform — three purpose-built experiences")

portals = [
    ("🛡  Admin Portal", NAVY, [
        "System-wide operations dashboard",
        "Tenant onboarding & management",
        "User management across all tenants",
        "Vendor approval workflow",
        "Platform health monitoring",
    ]),
    ("🏢  Business Portal", GREEN, [
        "Buyer-side document creation (PO, Invoice)",
        "Partner network management & tier control",
        "Real-time shipment tracking",
        "Vendor browsing & connection management",
        "Analytics: OTIF, invoice match, fill rate",
        "API keys & webhook configuration",
    ]),
    ("🏭  Vendor Portal", BLUE, [
        "Respond to POs, submit invoices",
        "Product catalog management",
        "Quotes & RFQ response workflow",
        "Payments & remittance tracking",
        "Returns & RMA management",
        "Compliance document vault",
        "Team management with role-based access",
        "Performance analytics (OTIF, KPIs)",
    ]),
]

col_w = Inches(4.1)
x_positions = [Inches(0.3), Inches(4.55), Inches(8.8)]

for i, (title, color, items) in enumerate(portals):
    x = x_positions[i]
    # Card
    add_rect(slide, x, Inches(1.45), col_w, Inches(5.7), fill_rgb=WHITE,
             line_rgb=BORDER, line_width=Pt(0.75))
    # Top accent
    add_rect(slide, x, Inches(1.45), col_w, Inches(0.08), fill_rgb=color)
    # Header band
    add_rect(slide, x, Inches(1.53), col_w, Inches(0.65),
             fill_rgb=RGBColor(min(255,color[0]//5+200), min(255,color[1]//5+200), min(255,color[2]//5+200)))
    add_text(slide, title, x + Inches(0.15), Inches(1.55), col_w - Inches(0.3), Inches(0.55),
             font_size=Pt(13), bold=True, color=color)

    tx = slide.shapes.add_textbox(x + Inches(0.15), Inches(2.25),
                                  col_w - Inches(0.3), Inches(4.7))
    tf = tx.text_frame
    tf.word_wrap = True
    for item in items:
        add_para(tf, f"▸  {item}", font_size=Pt(10.5), color=DARK_GRAY, space_before=Pt(6))


# ══════════════════════════════════════════════════════════════════
# SLIDE 5 — VENDOR PORTAL DEEP DIVE
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, Inches(0), Inches(0), prs.slide_width, prs.slide_height, fill_rgb=LIGHT_GRAY)
header_bar(slide, "Vendor Portal — Feature Deep Dive", "16 pages · complete supplier operations suite")

features = [
    ("Dashboard",      NAVY,    ["Greeting hero + OTIF KPIs", "Attention strips (alerts + OTIF gap)", "Activity feed + quick actions"]),
    ("Documents",      BLUE,    ["Receive & acknowledge POs", "Line-item confirmation with qty", "Accept / Reject with reason"]),
    ("Invoices",       GREEN,   ["Create & submit invoices", "Outstanding / Overdue / Paid summary", "Status tracking (DRAFT→PAID)"]),
    ("Shipments",      AMBER,   ["Track with expandable history", "Create ASN with carrier/tracking", "In Transit live indicators"]),
    ("Catalog",        PURPLE,  ["10 pre-loaded SKUs", "Category filter + search", "Add/Edit/Delete products"]),
    ("Quotes / RFQ",   BLUE,    ["Receive & respond to RFQs", "Line-level price + lead time editor", "Submit/Save draft workflow"]),
    ("Returns & RMA",  AMBER,   ["Pipeline: Request→Approved→Closed", "Issue credit notes", "Full line-item breakdown"]),
    ("Payments",       GREEN,   ["Remittance advice per payment", "SEPA/SWIFT/ACH/WIRE methods", "Early-pay discount tracking"]),
    ("Doc Vault",      NAVY,    ["ISO, Insurance, W-8, Audit certs", "Expiry countdown & alerts", "Upload with type/issuer/dates"]),
    ("Team",           PURPLE,  ["5 roles: Admin→Viewer", "Invite / Suspend / Remove", "MFA status per user"]),
    ("Labels",         BLUE,    ["5 types: Shipping/Pallet/HU/Box/Custom", "Live SVG preview + print", "SSCC-18 barcode support"]),
    ("Analytics",      GREEN,   ["OTIF rings, perfect order, fill rate", "Financial KPIs (DSO, Cash-to-Cash)", "6-month trend charts"]),
]

cols = 4
rows = 3
col_w = Inches(3.15)
row_h = Inches(1.75)
x0 = Inches(0.3)
y0 = Inches(1.45)

for idx, (name, color, pts) in enumerate(features):
    col = idx % cols
    row = idx // cols
    x = x0 + col * (col_w + Inches(0.07))
    y = y0 + row * (row_h + Inches(0.07))

    add_rect(slide, x, y, col_w, row_h, fill_rgb=WHITE, line_rgb=BORDER, line_width=Pt(0.75))
    add_rect(slide, x, y, col_w, Inches(0.07), fill_rgb=color)
    add_text(slide, name, x + Inches(0.12), y + Inches(0.1),
             col_w - Inches(0.24), Inches(0.32),
             font_size=Pt(10), bold=True, color=color)

    tx = slide.shapes.add_textbox(x + Inches(0.12), y + Inches(0.42),
                                  col_w - Inches(0.24), Inches(1.2))
    tf = tx.text_frame
    tf.word_wrap = True
    for pt in pts:
        add_para(tf, f"• {pt}", font_size=Pt(9), color=MID_GRAY, space_before=Pt(3))


# ══════════════════════════════════════════════════════════════════
# SLIDE 6 — API BACKEND MODULES
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, Inches(0), Inches(0), prs.slide_width, prs.slide_height, fill_rgb=LIGHT_GRAY)
header_bar(slide, "API Backend — NestJS 10 + Fastify", "Port 3001 · 15+ modules · production-ready")

modules = [
    ("auth",          "JWT RS256 · Argon2 passwords · TOTP MFA",                          GREEN),
    ("documents",     "Full lifecycle PO/INV/ASN · AES-256-GCM · audit trail",            BLUE),
    ("partners",      "Connection requests · 3-tier access (Standard/Preferred/Trusted)",  NAVY),
    ("vendors",       "Self-registration · admin approval · vendor profiles",               GREEN),
    ("carriers",      "FedEx/UPS/DHL adapters · label generation · rate lookup",           AMBER),
    ("tracking",      "Real-time shipment tracking · ETA calc · event enrichment",         BLUE),
    ("analytics",     "OTIF · document trends · invoice match rate · KPI dashboards",      PURPLE),
    ("api-keys",      "Scoped key generation · rate limiting · usage tracking",            NAVY),
    ("webhooks",      "Event subscriptions · HMAC-SHA256 signing · BullMQ delivery",       AMBER),
    ("notifications", "Email/SMS alerts · per-user preferences",                           GREEN),
    ("search",        "Elasticsearch full-text document search",                           BLUE),
    ("storage",       "MinIO S3 integration · encrypted file uploads",                     PURPLE),
    ("tenants",       "PostgreSQL RLS · multi-tenant context · plan management",           NAVY),
    ("users",         "Roles: SUPER_ADMIN/ADMIN/MANAGER/OPERATOR/VIEWER",                  GREEN),
    ("ai",            "Claude Anthropic invoice matching · document classification",       AMBER),
]

col_w = Inches(6.2)
y0 = Inches(1.45)
row_h = Inches(0.355)
mid = Inches(6.75)

for i, (mod, desc, color) in enumerate(modules):
    col = i % 2
    row = i // 2
    x = Inches(0.3) if col == 0 else mid
    y = y0 + row * row_h

    add_rect(slide, x, y, col_w, row_h - Inches(0.03), fill_rgb=WHITE,
             line_rgb=BORDER, line_width=Pt(0.5))
    add_rect(slide, x, y, Inches(0.06), row_h - Inches(0.03), fill_rgb=color)

    add_text(slide, mod, x + Inches(0.12), y + Inches(0.02),
             Inches(1.2), Inches(0.3),
             font_size=Pt(9.5), bold=True, color=color)
    add_text(slide, desc, x + Inches(1.35), y + Inches(0.02),
             col_w - Inches(1.5), Inches(0.3),
             font_size=Pt(9), color=MID_GRAY)


# ══════════════════════════════════════════════════════════════════
# SLIDE 7 — DOCUMENT LIFECYCLE
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, Inches(0), Inches(0), prs.slide_width, prs.slide_height, fill_rgb=LIGHT_GRAY)
header_bar(slide, "Document Lifecycle & Types", "End-to-end supply chain document exchange")

# Pipeline arrow
stages = [
    ("DRAFT",        DARK_GRAY, Inches(0.4)),
    ("SENT",         BLUE,      Inches(2.45)),
    ("ACKNOWLEDGED", AMBER,     Inches(4.5)),
    ("ACCEPTED",     GREEN,     Inches(7.2)),
    ("REJECTED",     RGBColor(0xBB, 0x00, 0x00), Inches(9.8)),
]
for label, color, x in stages:
    w = Inches(1.9)
    add_rect(slide, x, Inches(1.5), w, Inches(0.7), fill_rgb=color)
    add_text(slide, label, x, Inches(1.5), w, Inches(0.7),
             font_size=Pt(9.5), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if x.inches < 9.5:
        add_text(slide, "→", x + w, Inches(1.52), Inches(0.45), Inches(0.65),
                 font_size=Pt(18), bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# Document types grid
doc_types = [
    ("Purchase Order (PO)",       "Buyer → Vendor. Line items, delivery terms, pricing.",   NAVY),
    ("Invoice",                    "Vendor → Buyer. Payment request with 3-way match.",       GREEN),
    ("Advance Ship Notice (ASN)",  "Vendor → Buyer. Pre-shipment detail, SSCC barcodes.",     BLUE),
    ("Delivery Note",              "Confirms physical delivery with timestamps.",              AMBER),
    ("Bill of Lading (BOL)",       "Carrier-issued transport contract.",                      PURPLE),
    ("Credit Note",                "Issued against returns/RMAs to reverse charges.",         GREEN),
    ("Customs Declaration",        "Cross-border shipment regulatory documents.",             NAVY),
    ("Shipping Label",             "ZPL/PDF/PNG — FedEx, UPS, DHL, USPS, SSCC-18.",          BLUE),
]

y0 = Inches(2.5)
col_w = Inches(6.0)
row_h = Inches(0.62)
for i, (name, desc, color) in enumerate(doc_types):
    col = i % 2
    row = i // 2
    x = Inches(0.3) if col == 0 else Inches(6.7)
    y = y0 + row * row_h

    add_rect(slide, x, y, col_w, row_h - Inches(0.05), fill_rgb=WHITE,
             line_rgb=BORDER, line_width=Pt(0.75))
    add_rect(slide, x, y, Inches(0.07), row_h - Inches(0.05), fill_rgb=color)
    add_text(slide, name, x + Inches(0.15), y + Inches(0.04),
             Inches(2.2), Inches(0.3), font_size=Pt(10), bold=True, color=color)
    add_text(slide, desc, x + Inches(2.4), y + Inches(0.06),
             col_w - Inches(2.6), Inches(0.42), font_size=Pt(9.5), color=MID_GRAY)

# 3-way matching note
add_rect(slide, Inches(0.3), Inches(6.25), Inches(12.7), Inches(0.65),
         fill_rgb=RGBColor(0xEF, 0xF5, 0xFF), line_rgb=BLUE, line_width=Pt(1))
add_text(slide, "3-Way Matching:  Purchase Order  ↔  Invoice  ↔  Advance Ship Notice  "
         "— automated reconciliation with AI-assisted discrepancy detection via Claude Anthropic",
         Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.55),
         font_size=Pt(10), color=BLUE, bold=True)


# ══════════════════════════════════════════════════════════════════
# SLIDE 8 — SECURITY & COMPLIANCE
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, Inches(0), Inches(0), prs.slide_width, prs.slide_height, fill_rgb=LIGHT_GRAY)
header_bar(slide, "Security & Compliance", "Enterprise-grade security at every layer")

sec_items = [
    ("🔐  Authentication",  GREEN, [
        "JWT RS256 — asymmetric key signing",
        "Argon2id password hashing",
        "TOTP-based multi-factor authentication (MFA)",
        "Session management with Redis",
    ]),
    ("🛡  Data Encryption", BLUE, [
        "AES-256-GCM encryption for documents at rest",
        "RSA-PSS 2048-bit digital signatures per document",
        "TLS 1.3 for all data in transit",
        "IBAN/account data masked by default",
    ]),
    ("🏢  Multi-Tenancy",   NAVY, [
        "PostgreSQL Row-Level Security (RLS)",
        "app.tenant_id session variable injection",
        "Zero cross-tenant data leakage by design",
        "Tenant-scoped API keys with fine-grained scopes",
    ]),
    ("🔑  Access Control",  AMBER, [
        "5 platform roles: SUPER_ADMIN → VIEWER",
        "5 vendor roles: Admin/Manager/Finance/Logistics/Viewer",
        "Scoped API keys (po:read, invoice:write, etc.)",
        "Per-connection data-sharing tiers",
    ]),
    ("📋  Audit & Compliance", PURPLE, [
        "Full audit trail on every document state change",
        "Webhook HMAC-SHA256 signature verification",
        "W-9/W-8 form status tracking per vendor",
        "Compliance document vault with expiry alerts",
    ]),
    ("🌐  API Security",    GREEN, [
        "Kong API Gateway — rate limiting & request inspection",
        "API key prefix format: sf_live_*",
        "Input validation via Zod schemas on all endpoints",
        "CORS, helmet headers, injection prevention",
    ]),
]

col_w = Inches(4.0)
row_h = Inches(2.0)
x0, y0 = Inches(0.3), Inches(1.45)

for i, (title, color, pts) in enumerate(sec_items):
    col = i % 3
    row = i // 3
    x = x0 + col * (col_w + Inches(0.22))
    y = y0 + row * (row_h + Inches(0.08))

    add_rect(slide, x, y, col_w, row_h, fill_rgb=WHITE, line_rgb=BORDER, line_width=Pt(0.75))
    add_rect(slide, x, y, col_w, Inches(0.07), fill_rgb=color)
    add_text(slide, title, x + Inches(0.12), y + Inches(0.1),
             col_w - Inches(0.24), Inches(0.35), font_size=Pt(10.5), bold=True, color=color)
    tx = slide.shapes.add_textbox(x + Inches(0.12), y + Inches(0.5),
                                  col_w - Inches(0.24), Inches(1.4))
    tf = tx.text_frame
    tf.word_wrap = True
    for pt in pts:
        add_para(tf, f"• {pt}", font_size=Pt(9.5), color=MID_GRAY, space_before=Pt(4))


# ══════════════════════════════════════════════════════════════════
# SLIDE 9 — AI & INTEGRATIONS
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, Inches(0), Inches(0), prs.slide_width, prs.slide_height, fill_rgb=LIGHT_GRAY)
header_bar(slide, "AI & Integrations", "Claude Anthropic MCP · EDI · Webhooks · Data Feed")

# MCP block
add_rect(slide, Inches(0.3), Inches(1.45), Inches(6.1), Inches(5.7),
         fill_rgb=WHITE, line_rgb=BORDER, line_width=Pt(0.75))
add_rect(slide, Inches(0.3), Inches(1.45), Inches(6.1), Inches(0.07), fill_rgb=PURPLE)
add_text(slide, "🤖  MCP Server (Claude AI Integration)",
         Inches(0.45), Inches(1.55), Inches(5.8), Inches(0.4),
         font_size=Pt(12), bold=True, color=PURPLE)
add_text(slide, "8 tools exposed for agentic AI workflows",
         Inches(0.45), Inches(1.9), Inches(5.8), Inches(0.3),
         font_size=Pt(10), color=MID_GRAY, italic=True)

mcp_tools = [
    ("create_purchase_order",  "Generate PO with line items, delivery, payment terms"),
    ("get_document_status",    "Fetch document status, history, lifecycle events"),
    ("track_shipment",         "Real-time tracking with location and ETA"),
    ("match_invoice_to_po",    "AI 3-way matching — invoice ↔ PO ↔ ASN"),
    ("search_documents",       "Full-text search across documents with filters"),
    ("generate_label",         "Create shipping labels for FedEx/UPS/DHL"),
    ("vendor_performance",     "Get vendor KPIs — OTIF, quality, cost metrics"),
    ("approve_document",       "Approve/reject documents (PO/Invoice/ASN)"),
]
y_mcp = Inches(2.25)
for tool, desc in mcp_tools:
    add_rect(slide, Inches(0.45), y_mcp, Inches(5.8), Inches(0.52),
             fill_rgb=RGBColor(0xF5, 0xF0, 0xFF), line_rgb=RGBColor(0xD8, 0xC8, 0xF0), line_width=Pt(0.5))
    add_text(slide, tool, Inches(0.57), y_mcp + Inches(0.04),
             Inches(2.1), Inches(0.25), font_size=Pt(8.5), bold=True, color=PURPLE)
    add_text(slide, desc, Inches(2.7), y_mcp + Inches(0.06),
             Inches(3.3), Inches(0.35), font_size=Pt(8.5), color=MID_GRAY)
    y_mcp += Inches(0.57)

# Right column
right_x = Inches(6.7)

# EDI
add_rect(slide, right_x, Inches(1.45), Inches(6.3), Inches(2.0),
         fill_rgb=WHITE, line_rgb=BORDER, line_width=Pt(0.75))
add_rect(slide, right_x, Inches(1.45), Inches(6.3), Inches(0.07), fill_rgb=AMBER)
add_text(slide, "📡  EDI — @supplyforge/edi Package",
         right_x + Inches(0.15), Inches(1.55), Inches(6.0), Inches(0.35),
         font_size=Pt(11), bold=True, color=AMBER)
edi_items = [
    "X12 850 (Purchase Order)  ·  X12 810 (Invoice)  ·  X12 856 (ASN)  ·  X12 997 (Ack)",
    "EDIFACT ORDERS  ·  INVOIC  ·  DESADV  (Despatch Advice)",
    "Bidirectional: generate outbound + parse inbound EDI documents",
    "Enables direct legacy ERP integration (SAP, Oracle, etc.)",
]
tx = slide.shapes.add_textbox(right_x + Inches(0.15), Inches(1.95),
                              Inches(6.0), Inches(1.4))
tf = tx.text_frame; tf.word_wrap = True
for item in edi_items:
    add_para(tf, f"▸  {item}", font_size=Pt(9.5), color=MID_GRAY, space_before=Pt(5))

# Webhooks
add_rect(slide, right_x, Inches(3.55), Inches(6.3), Inches(1.6),
         fill_rgb=WHITE, line_rgb=BORDER, line_width=Pt(0.75))
add_rect(slide, right_x, Inches(3.55), Inches(6.3), Inches(0.07), fill_rgb=BLUE)
add_text(slide, "🔔  Webhook & Event System",
         right_x + Inches(0.15), Inches(3.65), Inches(6.0), Inches(0.35),
         font_size=Pt(11), bold=True, color=BLUE)
tx = slide.shapes.add_textbox(right_x + Inches(0.15), Inches(4.05),
                              Inches(6.0), Inches(1.0))
tf = tx.text_frame; tf.word_wrap = True
for item in ["HMAC-SHA256 signed payloads for security verification",
             "BullMQ async delivery with retry logic",
             "Events: document.created/sent/accepted/rejected, payment.sent, shipment.updated"]:
    add_para(tf, f"▸  {item}", font_size=Pt(9.5), color=MID_GRAY, space_before=Pt(5))

# Data Feed
add_rect(slide, right_x, Inches(5.25), Inches(6.3), Inches(1.85),
         fill_rgb=WHITE, line_rgb=BORDER, line_width=Pt(0.75))
add_rect(slide, right_x, Inches(5.25), Inches(6.3), Inches(0.07), fill_rgb=GREEN)
add_text(slide, "📊  Data Feed Subscriptions",
         right_x + Inches(0.15), Inches(5.35), Inches(6.0), Inches(0.35),
         font_size=Pt(11), bold=True, color=GREEN)
tx = slide.shapes.add_textbox(right_x + Inches(0.15), Inches(5.75),
                              Inches(6.0), Inches(1.2))
tf = tx.text_frame; tf.word_wrap = True
for item in ["Types: Inventory · Orders · Shipments · Pricing · Capacity",
             "Delivery: WEBHOOK · API_POLL · SFTP · AS2",
             "Vendor-controlled: choose what data to share with each buyer"]:
    add_para(tf, f"▸  {item}", font_size=Pt(9.5), color=MID_GRAY, space_before=Pt(5))


# ══════════════════════════════════════════════════════════════════
# SLIDE 10 — SHARED PACKAGES
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, Inches(0), Inches(0), prs.slide_width, prs.slide_height, fill_rgb=LIGHT_GRAY)
header_bar(slide, "Shared Packages", "Reusable libraries across API & Web apps")

packages = [
    ("@supplyforge/types", NAVY, "TypeScript Interface Library",
     "Shared type definitions consumed by all apps.",
     ["Tenant, User, Document, PO, Invoice, ASN, Shipment",
      "WebhookEvent, Notification, ApiKey, Partner",
      "Carrier, TrackingEvent, AnalyticsKPI",
      "Enums: DocumentStatus, TenantType, UserRole, CarrierType",
      "All types exported as strict TypeScript interfaces"]),
    ("@supplyforge/validators", BLUE, "Zod Validation Schemas",
     "Runtime + compile-time validation for all request/response shapes.",
     ["Registration & login schema validation",
      "Document creation / update validators",
      "API key scope validation",
      "Partner connection request schemas",
      "Address, banking, tax form validators"]),
    ("@supplyforge/crypto", GREEN, "Cryptography Utilities",
     "All cryptographic operations in one auditable package.",
     ["AES-256-GCM encrypt/decrypt for document content",
      "RSA-PSS 2048-bit sign + verify (document integrity)",
      "HMAC-SHA256 for webhook payload signing",
      "SHA-256 general-purpose hashing",
      "SSCC-18 barcode calculation (GS1 standard)"]),
    ("@supplyforge/edi", AMBER, "EDI Generator & Parser",
     "Bi-directional X12 and EDIFACT processing.",
     ["X12: 850 PO · 810 Invoice · 856 ASN · 997 Acknowledgement",
      "EDIFACT: ORDERS · INVOIC · DESADV",
      "Generate outbound EDI from internal document objects",
      "Parse inbound EDI into normalised TypeScript objects",
      "Enables SAP / Oracle / legacy ERP integration"]),
]

col_w = Inches(3.05)
x0 = Inches(0.28)
for i, (name, color, subtitle, summary, items) in enumerate(packages):
    x = x0 + i * (col_w + Inches(0.1))
    add_rect(slide, x, Inches(1.45), col_w, Inches(5.7),
             fill_rgb=WHITE, line_rgb=BORDER, line_width=Pt(0.75))
    add_rect(slide, x, Inches(1.45), col_w, Inches(0.07), fill_rgb=color)

    add_text(slide, name, x + Inches(0.12), Inches(1.55),
             col_w - Inches(0.24), Inches(0.38),
             font_size=Pt(10), bold=True, color=color)
    add_text(slide, subtitle, x + Inches(0.12), Inches(1.88),
             col_w - Inches(0.24), Inches(0.3),
             font_size=Pt(9), bold=True, color=NAVY)
    add_text(slide, summary, x + Inches(0.12), Inches(2.18),
             col_w - Inches(0.24), Inches(0.5),
             font_size=Pt(8.5), color=MID_GRAY, italic=True)

    # divider
    add_rect(slide, x + Inches(0.12), Inches(2.7),
             col_w - Inches(0.24), Inches(0.02), fill_rgb=BORDER)

    tx = slide.shapes.add_textbox(x + Inches(0.12), Inches(2.75),
                                  col_w - Inches(0.24), Inches(4.1))
    tf = tx.text_frame; tf.word_wrap = True
    for item in items:
        add_para(tf, f"▸  {item}", font_size=Pt(9), color=MID_GRAY, space_before=Pt(6))


# ══════════════════════════════════════════════════════════════════
# SLIDE 11 — INFRASTRUCTURE
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, Inches(0), Inches(0), prs.slide_width, prs.slide_height, fill_rgb=LIGHT_GRAY)
header_bar(slide, "Infrastructure & DevOps", "Docker · Kubernetes · Kong · Terraform")

infra = [
    ("PostgreSQL 16",    "5432",  NAVY,   "Primary database. Row-Level Security (RLS) for multi-tenancy.\napp.tenant_id injected per request for complete data isolation."),
    ("Redis 7",          "6379",  RGBColor(0xC0,0x00,0x00), "Session cache + rate limiting + BullMQ job queues.\nAsync document processing and webhook delivery queues."),
    ("RabbitMQ 3",       "5672 / 15672", AMBER, "Event message bus. Async document events, webhook delivery,\nEDI routing and partner notification fan-out."),
    ("MinIO (S3)",        "9000 / 9001", GREEN, "S3-compatible object storage for encrypted document files,\nshipping labels, attachments and audit logs."),
    ("Elasticsearch 8",  "9200",  BLUE,   "Full-text search across all documents, vendors, partners.\nIndexed by type, tenant, status and date range."),
    ("Kong Gateway",     "8000",  RGBColor(0x00,0x66,0xAA), "API gateway — rate limiting, JWT auth plugin,\nrequest/response transformation, access control lists."),
]

y0 = Inches(1.45)
row_h = Inches(0.88)
for i, (name, port, color, desc) in enumerate(infra):
    y = y0 + i * row_h
    add_rect(slide, Inches(0.3), y, Inches(12.7), row_h - Inches(0.06),
             fill_rgb=WHITE, line_rgb=BORDER, line_width=Pt(0.75))
    add_rect(slide, Inches(0.3), y, Inches(0.08), row_h - Inches(0.06), fill_rgb=color)

    # Service name
    add_text(slide, name, Inches(0.5), y + Inches(0.06),
             Inches(1.9), Inches(0.35), font_size=Pt(11), bold=True, color=color)
    # Port badge
    add_rect(slide, Inches(2.5), y + Inches(0.1), Inches(1.3), Inches(0.3),
             fill_rgb=RGBColor(color[0]//4+190, color[1]//4+190, color[2]//4+190),
             line_rgb=color, line_width=Pt(0.75))
    add_text(slide, port, Inches(2.5), y + Inches(0.1), Inches(1.3), Inches(0.3),
             font_size=Pt(8.5), bold=True, color=color, align=PP_ALIGN.CENTER)
    # Description
    add_text(slide, desc, Inches(4.0), y + Inches(0.05),
             Inches(8.8), Inches(0.75), font_size=Pt(9.5), color=MID_GRAY)

# Bottom row: deployment
add_rect(slide, Inches(0.3), Inches(6.9), Inches(4.0), Inches(0.45),
         fill_rgb=WHITE, line_rgb=BORDER, line_width=Pt(0.75))
add_text(slide, "🐳  Docker Compose — local dev environment",
         Inches(0.45), Inches(6.92), Inches(3.7), Inches(0.35),
         font_size=Pt(9.5), bold=True, color=NAVY)

add_rect(slide, Inches(4.5), Inches(6.9), Inches(4.0), Inches(0.45),
         fill_rgb=WHITE, line_rgb=BORDER, line_width=Pt(0.75))
add_text(slide, "☸️  Kubernetes — production deployments",
         Inches(4.65), Inches(6.92), Inches(3.7), Inches(0.35),
         font_size=Pt(9.5), bold=True, color=BLUE)

add_rect(slide, Inches(8.7), Inches(6.9), Inches(4.3), Inches(0.45),
         fill_rgb=WHITE, line_rgb=BORDER, line_width=Pt(0.75))
add_text(slide, "🌍  Terraform IaC — cloud infrastructure (AWS/GCP/Azure)",
         Inches(8.85), Inches(6.92), Inches(4.1), Inches(0.35),
         font_size=Pt(9.5), bold=True, color=GREEN)


# ══════════════════════════════════════════════════════════════════
# SLIDE 12 — KEY METRICS & ANALYTICS
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, Inches(0), Inches(0), prs.slide_width, prs.slide_height, fill_rgb=LIGHT_GRAY)
header_bar(slide, "Analytics & KPI Framework", "Real-time performance measurement across all dimensions")

kpi_sections = [
    ("📦  Supply Chain KPIs", GREEN, [
        ("OTIF",                  "On-Time In-Full delivery performance (target: 95%)"),
        ("Perfect Order Rate",    "Orders delivered on-time, in-full, damage-free, with correct docs"),
        ("Fill Rate",             "% of ordered quantity actually fulfilled"),
        ("On-Time Delivery",      "Shipments delivered by promised date"),
        ("Acknowledgement Compliance", "% of POs acknowledged within agreed window"),
    ]),
    ("💰  Financial KPIs", BLUE, [
        ("DSO",                   "Days Sales Outstanding — invoice to cash received"),
        ("Cash-to-Cash Cycle",    "Time from paying suppliers to receiving customer payment"),
        ("Invoice Match Rate",    "% of invoices matching PO with no discrepancies"),
        ("Invoice Accuracy",      "Error-free invoice rate"),
        ("Payment Terms Compliance", "% of payments made within agreed terms"),
        ("Early Pay Discount Rate", "Capture rate on available early-payment discounts"),
    ]),
    ("🏭  Procurement KPIs", AMBER, [
        ("PO Acknowledgement Time", "Average time to confirm receipt of purchase order"),
        ("Average Lead Time",       "Days from PO issuance to goods receipt"),
        ("Lead Time Compliance",    "% of deliveries within contracted lead time"),
        ("PO Acceptance Rate",      "% of POs accepted without modification"),
        ("Backorder Rate",          "% of orders with at least one line on backorder"),
    ]),
    ("🌱  Quality & ESG", PURPLE, [
        ("First Pass Yield",       "% of units passing QC on first inspection"),
        ("Quality Score",          "Composite quality rating per vendor"),
        ("Damage-Free Rate",       "% of shipments received without damage"),
        ("Carbon per Shipment",    "Scope 3 emissions tracking per delivery"),
        ("Supplier Score",         "Composite 0–100 vendor performance index"),
    ]),
]

col_w = Inches(3.05)
x0 = Inches(0.28)
for i, (title, color, kpis) in enumerate(kpi_sections):
    x = x0 + i * (col_w + Inches(0.1))
    add_rect(slide, x, Inches(1.45), col_w, Inches(5.7),
             fill_rgb=WHITE, line_rgb=BORDER, line_width=Pt(0.75))
    add_rect(slide, x, Inches(1.45), col_w, Inches(0.07), fill_rgb=color)
    add_text(slide, title, x + Inches(0.12), Inches(1.55),
             col_w - Inches(0.24), Inches(0.38),
             font_size=Pt(10.5), bold=True, color=color)

    tx = slide.shapes.add_textbox(x + Inches(0.12), Inches(2.0),
                                  col_w - Inches(0.24), Inches(5.0))
    tf = tx.text_frame; tf.word_wrap = True
    for kpi_name, kpi_desc in kpis:
        add_para(tf, kpi_name, font_size=Pt(10), bold=True, color=DARK_GRAY, space_before=Pt(8))
        add_para(tf, kpi_desc, font_size=Pt(9), color=MID_GRAY, space_before=Pt(2))


# ══════════════════════════════════════════════════════════════════
# SLIDE 13 — TECHNOLOGY STACK SUMMARY
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, Inches(0), Inches(0), prs.slide_width, prs.slide_height, fill_rgb=LIGHT_GRAY)
header_bar(slide, "Technology Stack", "Modern, production-proven technologies")

stack_items = [
    # (category, items, color, x, y, w, h)
    ("Frontend",      ["Next.js 14 (App Router)", "React 18", "TypeScript 5.4 (strict)", "TailwindCSS", "React Query (TanStack)", "Zod validation"],            BLUE,   Inches(0.3),  Inches(1.45), Inches(3.05), Inches(3.4)),
    ("Backend",       ["NestJS 10 + Fastify", "TypeScript 5.4", "Prisma ORM", "Class-validator / transformer", "Passport.js (JWT, API Key)", "BullMQ (Redis)"],  GREEN,  Inches(3.5),  Inches(1.45), Inches(3.05), Inches(3.4)),
    ("Data Stores",   ["PostgreSQL 16 (RLS)", "Redis 7", "Elasticsearch 8", "MinIO (S3-compat.)", "TimescaleDB-ready", "Prisma migrations"],                    NAVY,   Inches(6.7),  Inches(1.45), Inches(3.05), Inches(3.4)),
    ("Infra / DevOps",["Turborepo monorepo", "pnpm workspaces", "Docker Compose", "Kubernetes", "Kong API Gateway", "Terraform IaC"],                           AMBER,  Inches(9.9),  Inches(1.45), Inches(3.05), Inches(3.4)),
    ("Messaging",     ["RabbitMQ 3", "BullMQ job queues", "NestJS EventEmitter", "Webhook delivery", "HMAC-SHA256 signing"],                                     PURPLE, Inches(0.3),  Inches(5.0),  Inches(3.05), Inches(2.15)),
    ("Security",      ["JWT RS256", "Argon2id", "AES-256-GCM", "RSA-PSS 2048", "TOTP MFA", "TLS 1.3"],                                                          RGBColor(0xBB,0x00,0x00), Inches(3.5), Inches(5.0), Inches(3.05), Inches(2.15)),
    ("AI / Integr.",  ["Claude Anthropic (MCP)", "8 MCP tools", "X12 EDI", "EDIFACT", "SFTP · AS2", "MailHog (email dev)"],                                     BLUE,   Inches(6.7),  Inches(5.0),  Inches(3.05), Inches(2.15)),
    ("Language",      ["TypeScript 5.4 (strict mode)", "Full-stack type safety", "@supplyforge/types", "Shared validators", "Zod schemas", "100% typed APIs"],  GREEN,  Inches(9.9),  Inches(5.0),  Inches(3.05), Inches(2.15)),
]

for category, items, color, x, y, w, h in stack_items:
    add_rect(slide, x, y, w, h, fill_rgb=WHITE, line_rgb=BORDER, line_width=Pt(0.75))
    add_rect(slide, x, y, w, Inches(0.07), fill_rgb=color)
    add_text(slide, category, x + Inches(0.12), y + Inches(0.1),
             w - Inches(0.24), Inches(0.32), font_size=Pt(10), bold=True, color=color)
    tx = slide.shapes.add_textbox(x + Inches(0.12), y + Inches(0.48),
                                  w - Inches(0.24), h - Inches(0.6))
    tf = tx.text_frame; tf.word_wrap = True
    for item in items:
        add_para(tf, f"• {item}", font_size=Pt(9.5), color=MID_GRAY, space_before=Pt(4))


# ══════════════════════════════════════════════════════════════════
# SLIDE 14 — ROADMAP / THANK YOU
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, Inches(0), Inches(0), prs.slide_width, prs.slide_height, fill_rgb=DARK_NAVY)
add_rect(slide, Inches(0), Inches(3.8), prs.slide_width, Inches(0.04), fill_rgb=GREEN)

header_bar(slide, "Feature Roadmap", "Next development priorities")

roadmap = [
    ("Tier 2 Enhancements", AMBER, [
        "Dispute Center — formal dispute workflow with thread comms",
        "Contracts & Pricing Agreements — blanket POs, volume discounts",
        "Inventory Management — live stock levels per SKU/warehouse",
        "Messaging Center — threaded buyer-vendor communication",
        "Onboarding Wizard — step-by-step vendor setup flow",
        "Credit Notes — issue against invoices/returns",
    ]),
    ("Tier 3 Advanced", PURPLE, [
        "AI Assistant — Claude chat for 'Why is my OTIF low?'",
        "Demand Forecasting — rolling 3/6-month from PO history",
        "Early Payment Discounts — dynamic discounting dashboard",
        "Quality Management — NCRs, CAPAs, quality incidents",
        "EDI Configuration UI — self-service trading partner setup",
        "ESG / Sustainability — scope 3 emissions, diversity reporting",
    ]),
    ("Platform Enhancements", BLUE, [
        "API Keys self-service in vendor portal",
        "Webhook configuration UI for vendors",
        "Mobile PWA for on-the-go access",
        "Business portal analytics deep-dive",
        "Admin portal advanced reporting",
        "Multi-language (i18n) support",
    ]),
]

col_w = Inches(3.9)
x0 = Inches(0.35)
for i, (title, color, items) in enumerate(roadmap):
    x = x0 + i * (col_w + Inches(0.22))
    add_rect(slide, x, Inches(1.45), col_w, Inches(4.5),
             fill_rgb=RGBColor(0x1E, 0x34, 0x46), line_rgb=color, line_width=Pt(1.5))
    add_rect(slide, x, Inches(1.45), col_w, Inches(0.07), fill_rgb=color)
    add_text(slide, title, x + Inches(0.15), Inches(1.55),
             col_w - Inches(0.3), Inches(0.38),
             font_size=Pt(11), bold=True, color=color)
    tx = slide.shapes.add_textbox(x + Inches(0.15), Inches(2.0),
                                  col_w - Inches(0.3), Inches(3.8))
    tf = tx.text_frame; tf.word_wrap = True
    for item in items:
        add_para(tf, f"▸  {item}", font_size=Pt(10), color=WHITE, space_before=Pt(7))

# Closing message
add_text(slide, "SupplyForge — Built for the modern supply chain",
         Inches(1.0), Inches(6.15), Inches(11.3), Inches(0.55),
         font_size=Pt(16), bold=True, color=LIGHT_GREEN, align=PP_ALIGN.CENTER)
add_text(slide, "Multi-tenant · AI-powered · Enterprise-grade · Production-ready",
         Inches(1.0), Inches(6.65), Inches(11.3), Inches(0.4),
         font_size=Pt(11), color=MID_GRAY, italic=True, align=PP_ALIGN.CENTER)


# ── Save ─────────────────────────────────────────────────────────
out = r"c:\Users\Public\mydata\SupplyForge\SupplyForge_Summary.pptx"
prs.save(out)
print(f"✅  Saved: {out}")
print(f"📊  Slides: {len(prs.slides)}")
